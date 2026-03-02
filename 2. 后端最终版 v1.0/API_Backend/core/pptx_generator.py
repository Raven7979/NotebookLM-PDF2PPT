import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from typing import List, Dict, Any
from PIL import Image
from io import BytesIO

class PPTXGenerator:
    def __init__(self):
        # 13.33 x 7.5 inches (16:9)
        self.width = Inches(13.33)
        self.height = Inches(7.5)

    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        hex_color = hex_color.lstrip('#')
        if len(hex_color) == 6:
            return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
        return RGBColor(0, 0, 0)

    def generate(self, pages_data: List[Dict[str, Any]], output_path: str):
        """
        Generate PPTX from processed page data.
        pages_data item structure:
        {
            "image": PIL.Image (background),
            "text_blocks": List[Dict] (from Qwen)
        }
        """
        prs = Presentation()
        prs.slide_width = int(self.width)
        prs.slide_height = int(self.height)

        for page in pages_data:
            # Use blank layout (usually index 6 in standard template)
            slide_layout = prs.slide_layouts[6] 
            slide = prs.slides.add_slide(slide_layout)

            # 1. Add Background Image
            img = page['image']
            image_stream = BytesIO()
            img.save(image_stream, format="PNG")
            image_stream.seek(0)
            
            # Stretch to fit slide
            slide.shapes.add_picture(image_stream, 0, 0, width=self.width, height=self.height)

            # 2. Add Text Blocks (Overlay)
            text_blocks = page.get('text_blocks', [])
            for block in text_blocks:
                # Calculate position and size based on percentages
                # x, y, width, height are 0.0-1.0
                left = self.width * block.get('x', 0)
                top = self.height * block.get('y', 0)
                width = self.width * block.get('width', 0.1)
                height = self.height * block.get('height', 0.05)
                
                textbox = slide.shapes.add_textbox(left, top, width, height)
                tf = textbox.text_frame
                tf.word_wrap = True
                
                p = tf.paragraphs[0]
                p.text = block.get('text', '')
                
                # Apply styling
                if 'fontSize' in block:
                    # Swift code estimates font size in pixels, convert to Pt approx
                    # This might need tuning
                    p.font.size = Pt(block['fontSize'])
                
                if 'color' in block:
                    try:
                        p.font.color.rgb = self._hex_to_rgb(block['color'])
                    except:
                        pass
                
                if block.get('isBold', False):
                    p.font.bold = True

        prs.save(output_path)
